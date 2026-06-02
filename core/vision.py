"""이미지/영상 위주 문서 이해 — Ollama 멀티모달 모델(gemma)로 시각 콘텐츠 요약.

텍스트 추출이 빈약한 PDF/PPTX(슬라이드 이미지, 다이어그램, 사진 위주)를
시각 모델로 직접 읽어 요약 텍스트를 생성한다.
"""
import base64
import json
import os
import urllib.request

OLLAMA_URL   = "http://localhost:11434/api/chat"
VISION_MODEL = "gemma4:e2b"   # capabilities: vision, audio, thinking

# 한 번에 모델에 보낼 최대 이미지 수 (속도/메모리 균형)
MAX_PDF_PAGES = 6
MAX_PPTX_IMGS = 10


def _ollama_vision(images_b64: list[str], prompt: str, num_predict: int = 900) -> str:
    """이미지 목록 + 프롬프트를 멀티모달 모델에 전달. thinking 비활성화로 직접 답변 유도."""
    if not images_b64:
        return ""
    payload = json.dumps({
        "model": VISION_MODEL,
        "messages": [{"role": "user", "content": prompt, "images": images_b64}],
        "stream": False,
        "think": False,                      # thinking 토큰 소비 방지
        "options": {"temperature": 0.2, "num_predict": num_predict},
    }).encode()
    req = urllib.request.Request(
        OLLAMA_URL, data=payload,
        headers={"Content-Type": "application/json"}, method="POST"
    )
    with urllib.request.urlopen(req, timeout=300) as resp:
        body = json.loads(resp.read())
    return body.get("message", {}).get("content", "").strip()


def render_pdf_to_images(file_path: str, max_pages: int = MAX_PDF_PAGES, dpi: int = 110) -> list[str]:
    """PDF 페이지들을 PNG로 렌더 → base64 리스트"""
    import fitz
    imgs = []
    doc = fitz.open(file_path)
    n = min(len(doc), max_pages)
    for i in range(n):
        pix = doc[i].get_pixmap(dpi=dpi)
        imgs.append(base64.b64encode(pix.tobytes("png")).decode())
    doc.close()
    return imgs


def extract_pptx_images(file_path: str, max_imgs: int = MAX_PPTX_IMGS) -> list[str]:
    """PPTX 슬라이드에 임베디드된 그림(사진/다이어그램)을 base64로 추출"""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    imgs = []
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    imgs.append(base64.b64encode(blob).decode())
                except Exception:
                    continue
            if len(imgs) >= max_imgs:
                return imgs
    return imgs


def collect_images(file_path: str, source_type: str) -> list[str]:
    """파일 형식에 맞게 시각 콘텐츠(이미지)를 수집"""
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if source_type == "pdf" or ext == ".pdf":
            return render_pdf_to_images(file_path)
        if source_type == "pptx" or ext in (".pptx", ".ppt"):
            return extract_pptx_images(file_path)
        if ext in (".png", ".jpg", ".jpeg", ".webp", ".bmp"):
            with open(file_path, "rb") as f:
                return [base64.b64encode(f.read()).decode()]
    except Exception as e:
        print(f"[vision] 이미지 수집 실패 {file_path}: {e}")
    return []


VISION_PROMPT = """다음은 한 문서의 페이지/슬라이드 이미지들입니다.
이미지에 담긴 도표, 그림, 수식, 텍스트를 모두 살펴보고 이 문서가 무엇을 다루는지 한국어로 정리하세요.

다음 형식으로 작성하세요 (다른 군더더기 없이):
주제: (한 줄)
핵심 개념: (쉼표로 구분, 3~6개)
주요 내용:
- (불릿 3~6개)
요약: (3~4문장)

이미지에서 실제로 보이는 것만 근거로 작성하세요."""


def describe_document(file_path: str, source_type: str = "", title: str = "") -> str:
    """이미지 위주 문서를 시각 모델로 읽어 요약 텍스트(한국어)를 반환. 실패 시 빈 문자열."""
    images = collect_images(file_path, source_type)
    if not images:
        return ""
    prompt = VISION_PROMPT
    if title:
        prompt = f"문서 제목: {title}\n\n{prompt}"
    try:
        return _ollama_vision(images, prompt)
    except Exception as e:
        print(f"[vision] describe_document 실패 {file_path}: {e}")
        return ""


def is_sparse(content: str) -> bool:
    """텍스트가 빈약해 시각 이해가 필요한지 판단"""
    return len((content or "").strip()) < 200
