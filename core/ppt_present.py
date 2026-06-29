"""PPT 자동 발표 — 슬라이드 추출(텍스트+이미지) + 발표 스크립트 생성

PPTX는 PowerPoint COM으로, PDF는 PyMuPDF로 슬라이드(페이지)를 PNG로 내보내 발표 화면에 그대로 쓰고,
슬라이드별 텍스트/노트(+텍스트가 빈약하면 비전 모델로 이미지 설명)를 바탕으로
사용자의 말투 프로필을 반영한 발표 대본을 LLM으로 생성한다.

PDF는 PPTX를 PDF로 내보낸 자료가 많아(스피커 노트 없음) 본문 텍스트만으로 처리하고,
텍스트가 거의 없는 페이지(이미지 위주 슬라이드)는 렌더링된 페이지 이미지를 비전 모델로 보강한다.
"""
import os
from concurrent.futures import ThreadPoolExecutor
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from . import pattern, vision

SUPPORTED_EXTS = (".pptx", ".ppt", ".pdf")

SCRIPT_PROMPT = """당신은 발표자 본인입니다. 아래 슬라이드 내용을 바탕으로, 청중 앞에서 실제로 말하듯 자연스러운 발표 대본을 작성하세요.

{style_block}
{prev_block}
[슬라이드 {index}/{total}]
제목: {title}
본문:
{body}
{notes_block}
{vision_block}

규칙:
- 슬라이드에 적힌 문장을 그대로 읽지 말고, 청중에게 구어체로 설명하듯 풀어서 말하세요.
- 2~4문장, 한국어. 발표 대본 텍스트만 출력하세요(따옴표나 "대본:" 같은 군더더기 없이).
- 직전 슬라이드와 자연스럽게 이어지도록 하세요(같은 얘기를 반복하지 마세요)."""


def extract_slides(file_path: str) -> list[dict]:
    """슬라이드별 제목/본문/스피커노트 추출"""
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        title_text = ""
        body_parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.shape_type == 13:  # 제목 placeholder
                title_text = shape.text_frame.text.strip()
            else:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        body_parts.append(t)

        notes_text = ""
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        slides.append({
            "index": i + 1,
            "title": title_text,
            "body": "\n".join(body_parts),
            "notes": notes_text,
        })
    return slides


def export_slide_images(file_path: str, out_dir: str) -> list[str]:
    """PowerPoint COM으로 슬라이드 전체를 PNG로 내보내 파일명 목록(슬라이드 순서) 반환.
    PowerPoint가 설치돼 있어야 함. 실패 시 빈 리스트."""
    os.makedirs(out_dir, exist_ok=True)
    try:
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()
        try:
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            # WithWindow=False만으로는 일부 PowerPoint 버전에서 창이 잠깐 보임 — 영향 없음, 발표 종료 후 자동 종료됨
            pres = powerpoint.Presentations.Open(
                os.path.abspath(file_path), WithWindow=False
            )
            pres.Export(os.path.abspath(out_dir), "PNG")
            pres.Close()
        finally:
            try:
                powerpoint.Quit()
            except Exception:
                pass
            pythoncom.CoUninitialize()
    except Exception as e:
        print(f"[ppt_present] PowerPoint COM 내보내기 실패: {e}")
        return []

    # PowerPoint가 "Slide1.PNG", "Slide2.PNG" ... 형식으로 내보냄 (버전에 따라 0패딩 다름)
    files = [f for f in os.listdir(out_dir) if f.lower().endswith(".png")]

    def _slide_num(fname: str) -> int:
        digits = "".join(ch for ch in fname if ch.isdigit())
        return int(digits) if digits else 0

    files.sort(key=_slide_num)
    return files


def extract_slides_pdf(file_path: str) -> list[dict]:
    """PDF 페이지별 텍스트 추출 (제목/스피커노트 개념 없음 — body만 채움)"""
    import fitz
    slides = []
    doc = fitz.open(file_path)
    for i, page in enumerate(doc):
        text = page.get_text().strip()
        slides.append({
            "index": i + 1,
            "title": "",
            "body": text,
            "notes": "",
        })
    doc.close()
    return slides


def export_slide_images_pdf(file_path: str, out_dir: str, dpi: int = 150) -> list[str]:
    """PDF 페이지 전체를 PNG로 렌더 → 파일명 목록(페이지 순서) 반환"""
    import fitz
    os.makedirs(out_dir, exist_ok=True)
    files = []
    doc = fitz.open(file_path)
    for i, page in enumerate(doc):
        pix = page.get_pixmap(dpi=dpi)
        fname = f"page_{i + 1}.png"
        pix.save(os.path.join(out_dir, fname))
        files.append(fname)
    doc.close()
    return files


def _describe_image_for_script(image_path: str) -> str:
    """텍스트가 빈약한 슬라이드의 이미지를 비전 모델로 짧게 설명 (발표 대본용)"""
    try:
        with open(image_path, "rb") as f:
            import base64
            img_b64 = base64.b64encode(f.read()).decode()
        prompt = "이 슬라이드 이미지에 보이는 내용(도표/그림/텍스트)을 한국어 2~3문장으로 설명하세요. 군더더기 없이 내용만."
        return vision._ollama_vision([img_b64], prompt, num_predict=300)
    except Exception as e:
        print(f"[ppt_present] 슬라이드 이미지 설명 실패: {e}")
        return ""


def _style_block() -> str:
    """행동축 말투 프로필을 발표 대본 프롬프트에 주입"""
    from . import avatar as avatar_core
    profile = avatar_core.get_profile()

    def pval(key):
        return profile.get(key, {}).get("value", "") if profile else ""

    parts = []
    if pval("speech_style"):
        parts.append(f"- 말투: {pval('speech_style')}")
    if pval("persona"):
        parts.append(f"- 성격: {pval('persona')}")
    if pval("language_tone"):
        parts.append(f"- 톤: {pval('language_tone')}")
    if not parts:
        return ""
    return "발표자의 말투/성격(이 스타일을 일관되게 유지하세요):\n" + "\n".join(parts) + "\n"


def generate_script(slide: dict, total: int, style_block: str, prev_script: str, image_path: str | None) -> str:
    body = slide["body"] or "(본문 없음)"
    notes_block = f"\n스피커 노트:\n{slide['notes']}\n" if slide["notes"] else ""

    vision_block = ""
    is_sparse = len((slide["body"] or "") + (slide["notes"] or "")) < 60
    if is_sparse and image_path:
        desc = _describe_image_for_script(image_path)
        if desc:
            vision_block = f"\n슬라이드 이미지 설명:\n{desc}\n"

    prev_block = f"직전 슬라이드에서 한 말: \"{prev_script}\"\n" if prev_script else ""

    prompt = SCRIPT_PROMPT.format(
        style_block=style_block,
        prev_block=prev_block,
        index=slide["index"], total=total,
        title=slide["title"] or "(제목 없음)",
        body=body,
        notes_block=notes_block,
        vision_block=vision_block,
    )
    script = pattern._llm_call(prompt)
    return script.strip().strip('"')


def extract_and_export(file_path: str, out_dir: str) -> tuple[list[dict], list[str]]:
    """슬라이드 텍스트 추출 + 이미지 내보내기만 수행(대본 생성 전). regenerate_slide에서 재사용."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        slides = extract_slides_pdf(file_path)
        images = export_slide_images_pdf(file_path, out_dir)
    elif ext in (".pptx", ".ppt"):
        slides = extract_slides(file_path)
        images = export_slide_images(file_path, out_dir)
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {ext}")
    return slides, images


# 슬라이드당 LLM 호출을 몇 개씩 동시에 돌릴지. 로컬 Ollama+GPU를 다른 작업(XTTS/WebGL)과 공유하므로
# 너무 키우면 GPU 메모리 폭주로 드라이버가 리셋된 전례가 있어(server_restart_gotcha) 작게 유지.
_BATCH_SIZE = 2


def process_presentation(file_path: str, out_dir: str, progress_cb=None) -> list[dict]:
    """PPTX/PDF 파일 하나를 받아 슬라이드별 {index, image, script} 리스트로 변환.

    progress_cb(current, total)가 주어지면 슬라이드 추출 직후(current=0)와
    슬라이드별 대본 생성 완료마다 호출 — 오래 걸리는 LLM 호출 진행률 표시용.

    슬라이드는 _BATCH_SIZE개씩 묶어 동시에 생성한다(전부 직렬화하면 너무 느림).
    배치 내부 슬라이드들은 배치 시작 시점의 직전 대본을 공유해서 문맥을 잇고,
    배치가 끝나면 그 배치의 마지막 슬라이드 대본을 다음 배치의 prev_script로 넘긴다
    (완전 순차 대비 문맥 연결이 한 단계 느슨해지지만, 속도와 GPU 안전성을 위한 절충)."""
    slides, images = extract_and_export(file_path, out_dir)
    style_block = _style_block()
    total = len(slides)
    if progress_cb:
        progress_cb(0, total)

    def image_path_for(slide):
        image_name = images[slide["index"] - 1] if slide["index"] - 1 < len(images) else None
        return image_name, (os.path.join(out_dir, image_name) if image_name else None)

    result = []
    prev_script = ""
    done = 0
    for pos in range(0, total, _BATCH_SIZE):
        batch = slides[pos:pos + _BATCH_SIZE]
        batch_prev = prev_script
        with ThreadPoolExecutor(max_workers=len(batch)) as ex:
            futures = []
            for slide in batch:
                image_name, image_path = image_path_for(slide)
                futures.append((slide, image_name, ex.submit(generate_script, slide, total, style_block, batch_prev, image_path)))
            for slide, image_name, fut in futures:
                script = fut.result()
                result.append({
                    "index": slide["index"],
                    "title": slide["title"],
                    "image": image_name,
                    "script": script,
                })
                done += 1
                if progress_cb:
                    progress_cb(done, total)
        prev_script = result[-1]["script"]
    return result


def regenerate_slide(slide_index: int, raw_slides: list[dict], images: list[str], out_dir: str, prev_script: str) -> str:
    """저장된 원본 슬라이드 텍스트/이미지를 그대로 사용해 해당 슬라이드의 대본만 새로 생성."""
    slide = next((s for s in raw_slides if s["index"] == slide_index), None)
    if slide is None:
        raise ValueError(f"슬라이드 {slide_index}를 찾을 수 없습니다")
    image_name = images[slide_index - 1] if slide_index - 1 < len(images) else None
    image_path = os.path.join(out_dir, image_name) if image_name else None
    style_block = _style_block()
    return generate_script(slide, len(raw_slides), style_block, prev_script, image_path)
