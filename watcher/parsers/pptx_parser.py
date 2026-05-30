"""PPT(.pptx) 파서 — python-pptx 사용"""
import os
from pptx import Presentation


def parse_pptx(file_path: str) -> list[dict]:
    base_title = os.path.splitext(os.path.basename(file_path))[0]
    prs = Presentation(file_path)
    chunks = []

    slide_texts = []
    for i, slide in enumerate(prs.slides):
        title_text = ""
        body_parts = []
        notes_text = ""

        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.shape_type == 13:  # 제목
                    title_text = shape.text_frame.text.strip()
                else:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            body_parts.append(t)

        # 스피커 노트
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        slide_content = "\n".join(body_parts)
        if notes_text:
            slide_content += f"\n[노트] {notes_text}"

        if slide_content or title_text:
            slide_texts.append(f"[슬라이드 {i+1}] {title_text}\n{slide_content}")

    # 전체를 하나의 흐름으로 묶어 청킹
    full_text = "\n\n".join(slide_texts)
    CHUNK = 1500
    for idx in range(0, len(full_text), CHUNK):
        chunks.append({
            "title": f"{base_title} [{idx//CHUNK + 1}]" if len(full_text) > CHUNK else base_title,
            "content": full_text[idx:idx+CHUNK],
            "meta": {
                "source_type": "pptx",
                "file_path": file_path,
                "chunk_index": idx // CHUNK,
                "total_slides": len(prs.slides)
            }
        })

    return chunks if chunks else [{"title": base_title, "content": "", "meta": {"source_type": "pptx", "file_path": file_path}}]
